from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Couleurs ───────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x10, 0x45, 0x79)
MID_BLUE   = RGBColor(0x2D, 0x8B, 0xBA)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x2E)

# ── Dimensions (16:9 standard 13.33" × 7.5") ──────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────────────────

def add_oval(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(9, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill=None, line_color=None, line_w=Pt(1.5), radius=None):
    kind = 5 if radius is not None else 1
    shape = slide.shapes.add_shape(kind, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    if radius is not None and hasattr(shape, 'adjustments'):
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    return shape

def tb(slide, left, top, width, height, text, font="Amaranth", size=24,
        bold=False, color=BLACK, align=PP_ALIGN.LEFT, wrap=True, italic=False, spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.space_before = Pt(spacing)
    run = p.add_run()
    run.text       = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return box

def multiline_tb(slide, left, top, width, height, lines, font="Garet", size=17,
                  bold=False, color=BLACK, align=PP_ALIGN.LEFT, line_spacing=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, (txt, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.space_before = Pt(opts.get('space_before', line_spacing))
        run = p.add_run()
        run.text = txt
        run.font.name  = opts.get('font', font)
        run.font.size  = Pt(opts.get('size', size))
        run.font.bold  = opts.get('bold', bold)
        run.font.color.rgb = opts.get('color', color)
        run.font.italic = opts.get('italic', False)
    return box

def bg(slide):
    # White background rect
    add_rect(slide, 0, 0, W, H, fill=WHITE)
    # Top-left decorative oval
    add_oval(slide, -Inches(1.5), -Inches(1.5), Inches(5.5), Inches(5.5), LIGHT_BLUE)
    # Bottom-right decorative oval
    add_oval(slide, W - Inches(4), H - Inches(4), Inches(5.5), Inches(5.5), LIGHT_BLUE)

def border(slide):
    add_rect(slide, Inches(0.2), Inches(0.15), W - Inches(0.4), H - Inches(0.3),
             line_color=MID_BLUE, line_w=Pt(1.2), radius=0.03)

def slide_num(slide, num):
    tb(slide, W/2 - Inches(0.4), H - Inches(0.45), Inches(0.8), Inches(0.35),
       str(num).zfill(2), font="Garet", size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

def title_bar(slide, text):
    # Blue accent line under title
    add_rect(slide, Inches(0.4), Inches(1.25), Inches(2.5), Pt(3), fill=MID_BLUE)
    tb(slide, Inches(0.4), Inches(0.3), W - Inches(0.8), Inches(0.9),
       text, font="Amaranth", size=40, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)

def numbered_card(slide, num_str, label, desc, left, top, w, h):
    # Card background
    add_rect(slide, left, top, w, h, line_color=MID_BLUE, line_w=Pt(1), radius=0.05)
    # Number
    tb(slide, left + Inches(0.15), top + Inches(0.1), Inches(0.8), Inches(0.55),
       num_str, font="Garet", size=26, bold=True, color=DARK_BLUE)
    # Label
    tb(slide, left + Inches(0.15), top + Inches(0.55), w - Inches(0.3), Inches(0.5),
       label, font="Garet", size=14, bold=True, color=DARK_BLUE)
    # Description
    if desc:
        tb(slide, left + Inches(0.15), top + Inches(0.95), w - Inches(0.3), h - Inches(1.1),
           desc, font="Garet", size=12, color=BLACK, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COUVERTURE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)

# Diploma emoji top center
tb(s, W/2 - Inches(0.5), Inches(0.4), Inches(1), Inches(0.8),
   "🎓", font="Segoe UI Emoji", size=36, align=PP_ALIGN.CENTER)

# UNICONNECT
tb(s, Inches(0.5), Inches(1.2), W - Inches(1), Inches(1.8),
   "UNICONNECT", font="Amaranth", size=72, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Blue accent line
add_rect(s, W/2 - Inches(2.5), Inches(2.95), Inches(5), Pt(3), fill=MID_BLUE)

# Subtitle
tb(s, Inches(0.5), Inches(3.1), W - Inches(1), Inches(0.7),
   "Réseau Social Académique pour les Communautés Universitaires",
   font="Amaranth", size=22, bold=False, color=MID_BLUE, align=PP_ALIGN.CENTER)

# Separator dots
tb(s, Inches(0.5), Inches(3.9), W - Inches(1), Inches(0.4),
   "· · · · · · · · · · · · · · · · · · · · · · · · · · · · · · · · · · · · ·",
   font="Garet", size=11, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Name
tb(s, Inches(0.5), Inches(4.4), W - Inches(1), Inches(0.5),
   "[Votre Nom & Prénom]", font="Garet", size=18, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Project info
tb(s, Inches(0.5), Inches(5.0), W - Inches(1), Inches(0.5),
   "Projet Interdisciplinaire  |  Encadrant : M. Zidane",
   font="Garet", size=15, color=MID_BLUE, align=PP_ALIGN.CENTER)

slide_num(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "SOMMAIRE")

items_left = [
    ("01", "Introduction & Contexte"),
    ("02", "Problématique"),
    ("03", "Objectifs"),
    ("04", "Solution Proposée"),
    ("05", "Besoins Fonctionnels"),
]
items_right = [
    ("06", "Besoins Non Fonctionnels"),
    ("07", "Conception"),
    ("08", "Implémentation"),
    ("09", "Conclusion"),
]

col_w = Inches(5.5)
row_h = Inches(0.85)
start_top = Inches(1.5)
gap = Inches(0.15)

for i, (num, label) in enumerate(items_left):
    top = start_top + i * (row_h + gap)
    add_rect(s, Inches(0.5), top, col_w, row_h, line_color=MID_BLUE, line_w=Pt(1), radius=0.05)
    # Circle number
    add_oval(s, Inches(0.6), top + Inches(0.15), Inches(0.55), Inches(0.55), DARK_BLUE)
    tb(s, Inches(0.6), top + Inches(0.12), Inches(0.55), Inches(0.55),
       num, font="Garet", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, Inches(1.3), top + Inches(0.2), col_w - Inches(0.9), Inches(0.5),
       label, font="Garet", size=15, bold=False, color=DARK_BLUE)

for i, (num, label) in enumerate(items_right):
    top = start_top + i * (row_h + gap)
    add_rect(s, Inches(6.8), top, col_w, row_h, line_color=MID_BLUE, line_w=Pt(1), radius=0.05)
    add_oval(s, Inches(6.9), top + Inches(0.15), Inches(0.55), Inches(0.55), DARK_BLUE)
    tb(s, Inches(6.9), top + Inches(0.12), Inches(0.55), Inches(0.55),
       num, font="Garet", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, Inches(7.6), top + Inches(0.2), col_w - Inches(0.9), Inches(0.5),
       label, font="Garet", size=15, bold=False, color=DARK_BLUE)

slide_num(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION & CONTEXTE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "INTRODUCTION & CONTEXTE")

bullets = [
    "L'enseignement supérieur se digitalise — mais les étudiants communiquent encore via WhatsApp et Facebook",
    "Ces plateformes mélangent vie privée et vie académique",
    "Aucun espace numérique dédié à un campus n'existe aujourd'hui",
    "UniConnect est né de ce constat",
]
content_top = Inches(1.55)
for i, b in enumerate(bullets):
    t = content_top + i * Inches(1.3)
    add_oval(s, Inches(0.5), t + Inches(0.1), Inches(0.35), Inches(0.35), MID_BLUE)
    tb(s, Inches(0.55), t + Inches(0.05), Inches(0.35), Inches(0.35),
       "→", font="Garet", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, Inches(1.05), t, W - Inches(1.4), Inches(1.2),
       b, font="Garet", size=17, color=BLACK, wrap=True)

slide_num(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROBLÉMATIQUE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "PROBLÉMATIQUE")

problems = [
    ("P1", "Aucune plateforme académique dédiée"),
    ("P2", "Mélange vie privée & vie académique"),
    ("P3", "Ressources pédagogiques éparpillées"),
    ("P4", "Barrière linguistique : chacun à l'aise dans sa langue, exclusion involontaire"),
    ("P5", "Savoir perdu : compétences et solutions jamais partagées efficacement"),
]
cw = Inches(2.3)
ch = Inches(2.4)
start_l = Inches(0.45)
for i, (code, text) in enumerate(problems):
    l = start_l + i * (cw + Inches(0.2))
    t = Inches(1.55)
    add_rect(s, l, t, cw, ch, fill=DARK_BLUE if i == 0 else WHITE,
             line_color=MID_BLUE, line_w=Pt(1.2), radius=0.05)
    tb(s, l, t + Inches(0.15), cw, Inches(0.55),
       code, font="Garet", size=22, bold=True,
       color=WHITE if i == 0 else DARK_BLUE, align=PP_ALIGN.CENTER)
    add_rect(s, l + Inches(0.6), t + Inches(0.65), cw - Inches(1.2), Pt(2),
             fill=MID_BLUE if i > 0 else WHITE)
    tb(s, l + Inches(0.1), t + Inches(0.8), cw - Inches(0.2), ch - Inches(0.95),
       text, font="Garet", size=13,
       color=WHITE if i == 0 else BLACK,
       align=PP_ALIGN.CENTER, wrap=True)

slide_num(s, 4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — OBJECTIFS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "OBJECTIFS")

objectives = [
    ("01", "Espace académique centralisé"),
    ("02", "Faciliter l'échange et la collaboration"),
    ("03", "Centraliser les ressources pédagogiques"),
    ("04", "Accessibilité multilingue — AR / FR / EN"),
    ("05", "Environnement éthique et sain"),
    ("06", "Supporter plusieurs universités"),
]
cw = Inches(3.8)
ch = Inches(1.25)
gap_x = Inches(0.25)
gap_y = Inches(0.2)
sl = Inches(0.45)
st = Inches(1.55)
for i, (num, label) in enumerate(objectives):
    col = i % 3
    row = i // 3
    l = sl + col * (cw + gap_x)
    t = st + row * (ch + gap_y)
    add_rect(s, l, t, cw, ch, line_color=MID_BLUE, line_w=Pt(1.2), radius=0.05)
    add_oval(s, l + Inches(0.15), t + Inches(0.15), Inches(0.65), Inches(0.65), DARK_BLUE)
    tb(s, l + Inches(0.15), t + Inches(0.12), Inches(0.65), Inches(0.65),
       num, font="Garet", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, l + Inches(0.95), t + Inches(0.2), cw - Inches(1.1), ch - Inches(0.3),
       label, font="Garet", size=14, bold=False, color=DARK_BLUE, wrap=True)

slide_num(s, 5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SOLUTION PROPOSÉE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "SOLUTION PROPOSÉE")

sol_items = [
    ("🌐", "Application web full-stack dédiée à chaque université"),
    ("📰", "Fil d'actualité · Stories · Messagerie · Canaux · Bibliothèque · Notifications"),
    ("🛡️", "Panneau admin campus · Super-admin multi-universités"),
    ("🔒", "Isolation complète des données par campus"),
    ("✉️", "Inscription sécurisée par vérification d'email institutionnel\n— accès restreint au domaine @campus de chaque université"),
]
ct = Inches(1.55)
for i, (icon, text) in enumerate(sol_items):
    t = ct + i * Inches(1.08)
    add_rect(s, Inches(0.45), t, W - Inches(0.9), Inches(1.0),
             line_color=MID_BLUE if i % 2 == 0 else RGBColor(0xE8, 0xF4, 0xFB),
             fill=RGBColor(0xF5, 0xFA, 0xFF) if i % 2 != 0 else None,
             line_w=Pt(1), radius=0.04)
    tb(s, Inches(0.55), t + Inches(0.1), Inches(0.55), Inches(0.7),
       icon, font="Segoe UI Emoji", size=20, align=PP_ALIGN.CENTER)
    tb(s, Inches(1.2), t + Inches(0.1), W - Inches(1.7), Inches(0.8),
       text, font="Garet", size=15, color=BLACK, wrap=True)

slide_num(s, 6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BESOINS FONCTIONNELS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "BESOINS FONCTIONNELS")

roles = [
    {
        "icon": "🎓", "label": "Étudiant",
        "color": DARK_BLUE,
        "items": [
            "Inscription par email institutionnel vérifié",
            "Posts & stories · Like / Commentaire",
            "Messagerie directe et de groupe",
            "Suivre utilisateurs & canaux",
            "Avatar builder · Mode sombre/clair",
            "Interface trilingue (AR/FR/EN)",
        ]
    },
    {
        "icon": "⚙️", "label": "Admin Campus",
        "color": MID_BLUE,
        "items": [
            "Modérer posts & commentaires",
            "Bannir / suspendre des utilisateurs",
            "Gérer canaux & bibliothèque",
            "Suppression définitive (RGPD)",
        ]
    },
    {
        "icon": "👑", "label": "Super-Admin",
        "color": RGBColor(0x1A, 0x6B, 0x9A),
        "items": [
            "Gérer plusieurs universités",
            "Configurer les domaines email",
            "Attribuer les rôles admin",
        ]
    },
]

cw = Inches(3.9)
cl = Inches(0.45)
ct2 = Inches(1.5)
for i, role in enumerate(roles):
    l = cl + i * (cw + Inches(0.25))
    # Header card
    add_rect(s, l, ct2, cw, Inches(0.75), fill=role["color"], radius=0.05)
    tb(s, l, ct2 + Inches(0.1), Inches(0.7), Inches(0.55),
       role["icon"], font="Segoe UI Emoji", size=22, align=PP_ALIGN.CENTER)
    tb(s, l + Inches(0.65), ct2 + Inches(0.15), cw - Inches(0.7), Inches(0.5),
       role["label"], font="Garet", size=18, bold=True, color=WHITE)
    # Items
    for j, item in enumerate(role["items"]):
        it = ct2 + Inches(0.85) + j * Inches(0.82)
        add_rect(s, l, it, cw, Inches(0.75), line_color=MID_BLUE, line_w=Pt(0.8), radius=0.04)
        add_oval(s, l + Inches(0.1), it + Inches(0.28), Inches(0.18), Inches(0.18), MID_BLUE)
        tb(s, l + Inches(0.38), it + Inches(0.1), cw - Inches(0.45), Inches(0.6),
           item, font="Garet", size=12, color=BLACK, wrap=True)

slide_num(s, 7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BESOINS NON FONCTIONNELS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "BESOINS NON FONCTIONNELS")

nfr = [
    ("🔐", "Sécurité",        "Authentification · mots de passe chiffrés · contrôle d'accès par rôle"),
    ("⚡", "Performance",     "Pagination · build optimisé · chargement rapide"),
    ("🌍", "Accessibilité",   "3 langues · RTL auto · responsive · dark/light mode"),
    ("🧩", "Maintenabilité",  "MVC · 41 composants React · isolation par campus (university_id)"),
    ("🏢", "Multi-tenancy",   "Données complètement isolées entre universités"),
    ("📋", "Conformité RGPD", "Suppression définitive des données sur demande utilisateur"),
]

cw = Inches(3.9)
ch = Inches(1.35)
for i, (icon, label, desc) in enumerate(nfr):
    col = i % 3
    row = i // 3
    l = Inches(0.45) + col * (cw + Inches(0.25))
    t = Inches(1.55) + row * (ch + Inches(0.2))
    add_rect(s, l, t, cw, ch, line_color=MID_BLUE, line_w=Pt(1.2), radius=0.05)
    tb(s, l + Inches(0.1), t + Inches(0.1), Inches(0.5), Inches(0.5),
       icon, font="Segoe UI Emoji", size=20, align=PP_ALIGN.CENTER)
    tb(s, l + Inches(0.65), t + Inches(0.1), cw - Inches(0.75), Inches(0.45),
       label, font="Garet", size=14, bold=True, color=DARK_BLUE)
    add_rect(s, l + Inches(0.1), t + Inches(0.58), cw - Inches(0.2), Pt(1.5), fill=MID_BLUE)
    tb(s, l + Inches(0.1), t + Inches(0.65), cw - Inches(0.2), ch - Inches(0.75),
       desc, font="Garet", size=12, color=BLACK, wrap=True)

slide_num(s, 8)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ARCHITECTURE GLOBALE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "ARCHITECTURE GLOBALE")

layers = [
    ("Couche Présentation", "React · Inertia.js · Tailwind CSS · Vite", DARK_BLUE),
    ("Couche Métier",       "Laravel (Controllers · Middleware · Policies · Routes)", MID_BLUE),
    ("Couche Données",      "MySQL · Eloquent ORM · Migrations · Seeders", RGBColor(0x1A, 0x6B, 0x9A)),
]
lw = W - Inches(1.5)
lh = Inches(1.35)
ll = Inches(0.75)
lt = Inches(1.55)
arrow_w = Inches(0.4)
for i, (name, tech, color) in enumerate(layers):
    t = lt + i * (lh + Inches(0.25))
    add_rect(s, ll, t, lw, lh, fill=color, radius=0.04)
    tb(s, ll + Inches(0.2), t + Inches(0.1), lw * 0.4, Inches(0.55),
       name, font="Garet", size=18, bold=True, color=WHITE)
    tb(s, ll + Inches(0.2), t + Inches(0.65), lw - Inches(0.4), Inches(0.55),
       tech, font="Garet", size=14, color=RGBColor(0xD6, 0xEA, 0xF8))
    if i < 2:
        tb(s, ll + lw/2 - Inches(0.2), t + lh + Inches(0.05), Inches(0.4), Inches(0.2),
           "↕", font="Garet", size=18, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)

slide_num(s, 9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CAS D'UTILISATION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "CAS D'UTILISATION")

actors = [
    ("👤", "Visiteur",         "Consulter accueil\nEnvoyer message contact",                    WHITE, DARK_BLUE),
    ("🎓", "Étudiant",         "Posts · Stories · Messages\nSuivre · Canaux · Ressources",      DARK_BLUE, WHITE),
    ("⚙️", "Admin Campus",     "Modérer · Bannir\nGérer canaux & bibliothèque",                 MID_BLUE, WHITE),
    ("👑", "Super-Admin",      "Gérer universités\nConfigurer domaines · Rôles",                RGBColor(0x1A, 0x6B, 0x9A), WHITE),
]
cw = Inches(2.8)
ch = Inches(3.5)
cl2 = Inches(0.55)
for i, (icon, name, actions, bg_c, fg_c) in enumerate(actors):
    l = cl2 + i * (cw + Inches(0.3))
    t = Inches(1.55)
    add_rect(s, l, t, cw, ch, fill=bg_c if bg_c != WHITE else None,
             line_color=MID_BLUE, line_w=Pt(1.2), radius=0.05)
    tb(s, l, t + Inches(0.2), cw, Inches(0.7),
       icon, font="Segoe UI Emoji", size=32, align=PP_ALIGN.CENTER)
    tb(s, l, t + Inches(0.85), cw, Inches(0.55),
       name, font="Garet", size=17, bold=True, color=fg_c, align=PP_ALIGN.CENTER)
    add_rect(s, l + Inches(0.3), t + Inches(1.4), cw - Inches(0.6), Pt(1.5), fill=fg_c)
    lines = actions.split('\n')
    for j, line in enumerate(lines):
        tb(s, l + Inches(0.15), t + Inches(1.55) + j * Inches(0.75), cw - Inches(0.3), Inches(0.7),
           "• " + line, font="Garet", size=12, color=fg_c, align=PP_ALIGN.CENTER, wrap=True)

slide_num(s, 10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DIAGRAMMES DE SÉQUENCE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "DIAGRAMMES DE SÉQUENCE")

seqs = [
    ("📝", "Publier un post",
     "Soumettre contenu → Insérer post\n→ Enregistrer médias\n→ Mettre à jour karma"),
    ("💬", "Envoyer un message",
     "Envoyer → Vérifier conversation\n→ Créer si inexistante\n→ Insérer message + Notifier"),
    ("❤️", "Liker un post",
     "Cliquer Like → Vérifier campus\n→ Ajouter / retirer like\n→ Notifier l'auteur"),
]
cw = Inches(3.9)
ch = Inches(4.5)
for i, (icon, title, steps) in enumerate(seqs):
    l = Inches(0.45) + i * (cw + Inches(0.25))
    t = Inches(1.55)
    add_rect(s, l, t, cw, ch, line_color=MID_BLUE, line_w=Pt(1.2), radius=0.05)
    add_rect(s, l, t, cw, Inches(0.7), fill=DARK_BLUE, radius=0.05)
    tb(s, l, t + Inches(0.1), Inches(0.6), Inches(0.55),
       icon, font="Segoe UI Emoji", size=22, align=PP_ALIGN.CENTER)
    tb(s, l + Inches(0.55), t + Inches(0.12), cw - Inches(0.65), Inches(0.5),
       title, font="Garet", size=15, bold=True, color=WHITE)
    for j, step in enumerate(steps.split('\n')):
        st2 = t + Inches(0.85) + j * Inches(0.9)
        add_rect(s, l + Inches(0.2), st2, cw - Inches(0.4), Inches(0.75),
                 fill=LIGHT_BLUE, radius=0.04)
        tb(s, l + Inches(0.3), st2 + Inches(0.12), cw - Inches(0.5), Inches(0.55),
           step, font="Garet", size=12, color=DARK_BLUE, wrap=True)

slide_num(s, 11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DIAGRAMME DE CLASSES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "DIAGRAMME DE CLASSES")

add_rect(s, Inches(0.45), Inches(1.5), W - Inches(0.9), Inches(4.0),
         line_color=MID_BLUE, line_w=Pt(1), radius=0.04)

tb(s, Inches(0.55), Inches(1.6), W - Inches(1.1), Inches(0.5),
   "14 entités  ·  6 domaines fonctionnels",
   font="Garet", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

domains = [
    ("🔑", "Identité & Accès",    "University · User"),
    ("📰", "Feed Social",         "Post · PostMedia · PostItem\nLike · Save · Comment"),
    ("📸", "Stories",             "Story · StoryView · Reaction"),
    ("💬", "Messagerie",          "Conversation · Message"),
    ("👥", "Réseau Social",       "Follow · ChannelFollow · Channel"),
    ("📚", "Contenu Académique",  "Resource · Notification · ContactMessage"),
]
dw = Inches(3.85)
dh = Inches(1.4)
for i, (icon, name, entities) in enumerate(domains):
    col = i % 3
    row = i // 3
    l = Inches(0.55) + col * (dw + Inches(0.2))
    t = Inches(2.2) + row * (dh + Inches(0.15))
    add_rect(s, l, t, dw, dh, line_color=MID_BLUE, line_w=Pt(0.8), radius=0.04)
    tb(s, l + Inches(0.1), t + Inches(0.08), Inches(0.4), Inches(0.45),
       icon, font="Segoe UI Emoji", size=18)
    tb(s, l + Inches(0.55), t + Inches(0.08), dw - Inches(0.65), Inches(0.45),
       name, font="Garet", size=13, bold=True, color=DARK_BLUE)
    add_rect(s, l + Inches(0.1), t + Inches(0.52), dw - Inches(0.2), Pt(1), fill=LIGHT_BLUE)
    tb(s, l + Inches(0.1), t + Inches(0.6), dw - Inches(0.2), dh - Inches(0.7),
       entities, font="Garet", size=11, color=BLACK, wrap=True)

slide_num(s, 12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — STACK TECHNIQUE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "STACK TECHNIQUE")

stack = [
    ("⚙️", "Backend",       "Laravel  ·  MySQL\n(XAMPP en dev · Hostinger en production)",  DARK_BLUE),
    ("🎨", "Frontend",      "React  ·  Tailwind CSS",                                       MID_BLUE),
    ("🔗", "Lien & Build",  "Inertia.js  ·  Vite",                                          RGBColor(0x1A, 0x6B, 0x9A)),
]
sw = Inches(3.6)
sh = Inches(3.8)
sl2 = Inches(0.65)
for i, (icon, label, tech, color) in enumerate(stack):
    l = sl2 + i * (sw + Inches(0.45))
    t = Inches(1.6)
    add_rect(s, l, t, sw, sh, fill=color, radius=0.06)
    tb(s, l, t + Inches(0.3), sw, Inches(0.9),
       icon, font="Segoe UI Emoji", size=40, align=PP_ALIGN.CENTER)
    tb(s, l, t + Inches(1.2), sw, Inches(0.6),
       label, font="Garet", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, l + Inches(0.6), t + Inches(1.85), sw - Inches(1.2), Pt(1.5), fill=WHITE)
    for j, line in enumerate(tech.split('\n')):
        tb(s, l + Inches(0.1), t + Inches(2.0) + j * Inches(0.75), sw - Inches(0.2), Inches(0.65),
           line, font="Garet", size=15, color=RGBColor(0xD6, 0xEA, 0xF8),
           align=PP_ALIGN.CENTER, wrap=True)

slide_num(s, 13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — INTERFACE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "INTERFACE & DÉMONSTRATION")

screens = [
    ("📰", "Fil d'actualité",   "Posts · Stories · Navigation"),
    ("💬", "Messagerie",        "Direct · Groupes · Médias"),
    ("⚙️", "Panneau Admin",     "Modération · Utilisateurs · Canaux"),
    ("👑", "Panneau Super-Admin","Universités · Domaines · Rôles"),
]
sw2 = Inches(2.9)
sh2 = Inches(3.8)
sl3 = Inches(0.45)
for i, (icon, name, desc) in enumerate(screens):
    l = sl3 + i * (sw2 + Inches(0.2))
    t = Inches(1.6)
    add_rect(s, l, t, sw2, sh2, fill=LIGHT_BLUE, line_color=MID_BLUE, line_w=Pt(1.2), radius=0.05)
    # Placeholder for screenshot
    add_rect(s, l + Inches(0.1), t + Inches(0.1), sw2 - Inches(0.2), Inches(2.3),
             fill=RGBColor(0xC5, 0xDC, 0xEC), line_color=MID_BLUE, line_w=Pt(0.5), radius=0.03)
    tb(s, l + Inches(0.1), t + Inches(0.85), sw2 - Inches(0.2), Inches(0.6),
       icon, font="Segoe UI Emoji", size=30, align=PP_ALIGN.CENTER)
    tb(s, l + Inches(0.1), t + Inches(1.45), sw2 - Inches(0.2), Inches(0.45),
       "[ Capture d'écran ]", font="Garet", size=11, color=MID_BLUE, align=PP_ALIGN.CENTER, italic=True)
    add_rect(s, l + Inches(0.3), t + Inches(2.45), sw2 - Inches(0.6), Pt(1.5), fill=DARK_BLUE)
    tb(s, l, t + Inches(2.55), sw2, Inches(0.5),
       name, font="Garet", size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    tb(s, l, t + Inches(3.05), sw2, Inches(0.55),
       desc, font="Garet", size=11, color=MID_BLUE, align=PP_ALIGN.CENTER)

slide_num(s, 14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)
title_bar(s, "CONCLUSION")

# Block 1 - Ce que nous avons réalisé
add_rect(s, Inches(0.45), Inches(1.55), Inches(5.8), Inches(4.7),
         fill=DARK_BLUE, radius=0.05)
tb(s, Inches(0.55), Inches(1.65), Inches(5.6), Inches(0.55),
   "✅  Ce que nous avons réalisé", font="Garet", size=16, bold=True, color=WHITE)
add_rect(s, Inches(0.65), Inches(2.2), Inches(4.5), Pt(1.5), fill=MID_BLUE)
achieved = [
    "Plateforme académique complète",
    "9 fonctionnalités sociales",
    "Architecture multi-tenant",
    "Interface trilingue (AR/FR/EN)",
    "Environnement éthique · Conforme RGPD",
    "Inscription par email institutionnel",
]
for j, item in enumerate(achieved):
    tb(s, Inches(0.65), Inches(2.3) + j * Inches(0.58), Inches(5.5), Inches(0.52),
       "→  " + item, font="Garet", size=13, color=WHITE)

# Block 2 - Perspectives
add_rect(s, Inches(6.9), Inches(1.55), Inches(6.0), Inches(4.7),
         fill=MID_BLUE, radius=0.05)
tb(s, Inches(7.0), Inches(1.65), Inches(5.8), Inches(0.55),
   "🚀  Perspectives", font="Garet", size=16, bold=True, color=WHITE)
add_rect(s, Inches(7.1), Inches(2.2), Inches(4.5), Pt(1.5), fill=WHITE)
perspectives = [
    "WebSockets natifs (messagerie en polling actuellement)",
    "Application mobile (React Native / Flutter)",
    "IA : recommandations de contenu personnalisées",
    "Déploiement cloud (Hostinger)",
]
for j, item in enumerate(perspectives):
    tb(s, Inches(7.1), Inches(2.3) + j * Inches(0.72), Inches(5.7), Inches(0.65),
       "→  " + item, font="Garet", size=13, color=WHITE, wrap=True)

slide_num(s, 15)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — MERCI
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
border(s)

tb(s, Inches(0.5), Inches(1.0), W - Inches(1), Inches(1.0),
   "🎓", font="Segoe UI Emoji", size=48, align=PP_ALIGN.CENTER)
tb(s, Inches(0.5), Inches(2.1), W - Inches(1), Inches(1.4),
   "Merci pour votre attention", font="Amaranth", size=52, bold=True,
   color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_rect(s, W/2 - Inches(2.5), Inches(3.5), Inches(5), Pt(3), fill=MID_BLUE)
tb(s, Inches(0.5), Inches(3.65), W - Inches(1), Inches(0.7),
   "Nous sommes disponibles pour répondre à vos questions",
   font="Garet", size=18, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_rect(s, Inches(2.5), Inches(4.5), W - Inches(5), Inches(1.2),
         line_color=MID_BLUE, line_w=Pt(1), radius=0.05)
tb(s, Inches(2.5), Inches(4.65), W - Inches(5), Inches(0.5),
   "[Votre Nom]  ·  Encadrant : M. Zidane",
   font="Garet", size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

slide_num(s, 16)

# ── Save ───────────────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.expanduser("~"), "Desktop", "UniConnect_Presentation.pptx")
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f"Saved: {out}")
print(f"    Slides: {len(prs.slides)}")

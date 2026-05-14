from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Color palette ───────────────────────────────────────────────────────────
C_PRIMARY   = RGBColor(0x1A, 0x2E, 0x5A)   # deep navy
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # teal / cyan
C_ORANGE    = RGBColor(0xF4, 0xA2, 0x61)   # warm orange
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xF0, 0xF4, 0xF8)   # very light blue-gray
C_DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
C_GRAY      = RGBColor(0x64, 0x74, 0x8B)
C_GREEN     = RGBColor(0x10, 0xB9, 0x81)
C_RED       = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # blank layout

# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, border=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=C_DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=16, bold=False, color=C_DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def slide_header(slide, title, subtitle=None, accent=C_ACCENT):
    """Top banner with title."""
    add_rect(slide, 0, 0, 13.33, 1.4, fill=C_PRIMARY)
    add_rect(slide, 0, 1.35, 13.33, 0.08, fill=accent)
    add_text(slide, title, 0.5, 0.1, 11, 0.9,
             size=30, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.9, 11, 0.45,
                 size=14, color=accent, align=PP_ALIGN.LEFT)

def bullet_box(slide, items, l, t, w, h, title=None,
               dot_color=C_ACCENT, text_size=15, title_color=C_PRIMARY):
    """Render a list of bullet points inside a box."""
    if title:
        add_text(slide, title, l, t, w, 0.4,
                 size=16, bold=True, color=title_color)
        t += 0.42
        h -= 0.42
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        r1 = p.add_run()
        r1.text = "▸  "
        r1.font.color.rgb = dot_color
        r1.font.size = Pt(text_size)
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(text_size)
        r2.font.color.rgb = C_DARK_TEXT

def card(slide, l, t, w, h, bg=C_LIGHT, border=None):
    """Rounded-ish card (filled rectangle)."""
    s = add_rect(slide, l, t, w, h, fill=bg, border=border)
    return s

def footer(slide, page_num, total=19):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill=C_PRIMARY)
    add_text(slide, "UniConnect  |  Projet de Fin d'Études  |  2025–2026",
             0.3, 7.16, 10, 0.3, size=10, color=C_ACCENT)
    add_text(slide, f"{page_num} / {total}", 12.3, 7.16, 0.9, 0.3,
             size=10, color=C_WHITE, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

add_rect(s, 0, 0, 13.33, 7.5, fill=C_PRIMARY)
# diagonal accent strip
strip = s.shapes.add_shape(1, Inches(8), Inches(0), Inches(5.33), Inches(7.5))
strip.fill.solid(); strip.fill.fore_color.rgb = RGBColor(0x16, 0x27, 0x4E)
strip.line.fill.background()

add_rect(s, 0, 3.0, 13.33, 0.07, fill=C_ACCENT)

add_text(s, "🎓 UniConnect", 0.7, 0.8, 9, 1.1,
         size=52, bold=True, color=C_WHITE)
add_text(s, "Réseau Social Académique pour les Campus Universitaires",
         0.7, 1.9, 8.5, 0.8, size=20, color=C_ACCENT)
add_rect(s, 0.7, 3.2, 2.5, 0.07, fill=C_ORANGE)

add_text(s, "Projet de 3ème Année  |  Licence Informatique",
         0.7, 3.4, 8, 0.5, size=16, color=RGBColor(0xCB, 0xD5, 0xE1))
add_text(s, "Année universitaire 2025 – 2026",
         0.7, 3.9, 8, 0.5, size=14, color=C_GRAY)

# tech badges
for i, (label, bg) in enumerate([
    ("Laravel 12", RGBColor(0xFF, 0x29, 0x1D)),
    ("React 18",   RGBColor(0x61, 0xDA, 0xFB)),
    ("Inertia.js", RGBColor(0x92, 0x58, 0xFA)),
    ("Tailwind",   RGBColor(0x38, 0xBD, 0xF8)),
]):
    bx = 0.7 + i * 2.2
    add_rect(s, bx, 5.6, 1.9, 0.45, fill=bg)
    add_text(s, label, bx, 5.6, 1.9, 0.45,
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

footer(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Sommaire
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Plan de la Présentation")

sections = [
    ("01", "Introduction & Contexte"),
    ("02", "Problématique"),
    ("03", "Objectifs"),
    ("04", "Solution Proposée"),
    ("05", "Besoins Fonctionnels"),
    ("06", "Besoins Non Fonctionnels"),
    ("07", "Conception"),
    ("08", "Implémentation"),
    ("09", "Conclusion & Perspectives"),
]
cols = [0.4, 6.9]
for i, (num, title) in enumerate(sections):
    col = i % 2
    row = i // 2
    x = cols[col]
    y = 1.7 + row * 1.15
    card(s, x, y, 6.1, 0.9, bg=C_WHITE, border=C_ACCENT)
    add_text(s, num, x + 0.1, y + 0.08, 0.7, 0.7,
             size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, title, x + 0.85, y + 0.18, 5.0, 0.55,
             size=16, bold=True, color=C_PRIMARY)

footer(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Introduction / Contexte
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Introduction & Contexte",
             "Le numérique au service de la vie étudiante")

# Left column — context text
card(s, 0.4, 1.6, 7.5, 5.0, bg=C_WHITE, border=None)
add_text(s, "Contexte général", 0.6, 1.7, 7.0, 0.4,
         size=17, bold=True, color=C_PRIMARY)

items = [
    "L'enseignement supérieur se digitalise rapidement.",
    "Les étudiants manquent d'outils académiques adaptés à leurs campus.",
    "Les réseaux sociaux grand public (Facebook, Instagram…) mélangent vie privée et vie académique.",
    "Il n'existe pas de plateforme centralisée par université pour partager des ressources, collaborer et communiquer.",
    "UniConnect répond à ce vide en proposant un espace numérique dédié, éthique et multi-campus.",
]
bullet_box(s, items, 0.6, 2.15, 7.2, 4.2, text_size=14)

# Right column — stats cards
for i, (val, label, color) in enumerate([
    ("3 langues", "Arabe · Français · Anglais", C_ACCENT),
    ("Multi-campus", "Isolation des données par université", C_GREEN),
    ("Éthique", "Mode focus, karma, anonymat", C_ORANGE),
]):
    cy = 1.6 + i * 1.65
    card(s, 8.2, cy, 4.7, 1.45, bg=color)
    add_text(s, val,   8.35, cy + 0.1,  4.3, 0.6,
             size=22, bold=True, color=C_WHITE)
    add_text(s, label, 8.35, cy + 0.65, 4.3, 0.6,
             size=13, color=C_WHITE)

footer(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Problématique
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Problématique",
             "Quels problèmes cherche-t-on à résoudre ?")

problems = [
    ("📭", "Absence de plateforme\nacadémique dédiée",
     "Les étudiants n'ont aucun espace numérique propre à leur campus pour échanger sur des sujets académiques."),
    ("🔒", "Mélange vie privée\n& vie académique",
     "Les réseaux sociaux classiques ne distinguent pas les contenus personnels des contenus pédagogiques."),
    ("📦", "Ressources éparpillées",
     "Les documents de cours, liens utiles et annonces sont disséminés sur e-mails, groupes WhatsApp et Drive."),
    ("🌐", "Barrière linguistique",
     "Dans les universités maghrébines, la coexistence du français, de l'arabe et du darija freine la communication."),
]

for i, (icon, title, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.6 + row * 2.6
    card(s, x, y, 6.0, 2.3, bg=C_WHITE, border=C_RED)
    add_text(s, icon, x + 0.15, y + 0.15, 0.8, 0.8, size=30)
    add_text(s, title, x + 1.0, y + 0.1, 4.7, 0.9,
             size=16, bold=True, color=C_PRIMARY)
    add_text(s, desc,  x + 0.2, y + 1.0, 5.6, 1.1,
             size=13, color=C_GRAY)

footer(s, 4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Objectifs
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Objectifs du Projet",
             "Ce que UniConnect cherche à accomplir")

objectives = [
    ("🎯", "Créer un espace académique centralisé",
     "Un réseau social dédié à chaque université, isolé et sécurisé."),
    ("🤝", "Favoriser la collaboration entre pairs",
     "Posts, stories, canaux thématiques et messagerie directe."),
    ("📚", "Centraliser les ressources pédagogiques",
     "Bibliothèque de liens et documents filtrés par filière."),
    ("🌍", "Assurer l'accessibilité multilingue",
     "Support natif AR / FR / EN avec RTL automatique."),
    ("🛡️", "Garantir un environnement éthique",
     "Mode focus nocturne, filtrage de mots, karma, anonymat."),
    ("🏫", "Supporter plusieurs universités",
     "Architecture multi-tenant avec isolation complète des données."),
]

for i, (icon, title, desc) in enumerate(objectives):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.35
    y = 1.6 + row * 2.65
    card(s, x, y, 4.1, 2.4, bg=C_WHITE)
    add_rect(s, x, y, 4.1, 0.08, fill=C_ACCENT)
    add_text(s, icon, x + 0.15, y + 0.2, 0.7, 0.7, size=26)
    add_text(s, title, x + 0.2, y + 0.85, 3.7, 0.75,
             size=14, bold=True, color=C_PRIMARY)
    add_text(s, desc, x + 0.2, y + 1.55, 3.7, 0.7,
             size=12, color=C_GRAY)

footer(s, 5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Solution proposée
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Solution Proposée — UniConnect",
             "Un réseau social académique complet, éthique et multi-campus")

add_text(s,
    "UniConnect est une plateforme web full-stack offrant à chaque université son propre espace numérique, "
    "avec des fonctionnalités sociales riches (fil d'actualité, stories, messagerie, canaux) et des outils "
    "de collaboration académique (bibliothèque, profils, notifications) — le tout dans un cadre éthique et multilingue.",
    0.4, 1.55, 12.5, 1.0, size=14, color=C_DARK_TEXT)

features = [
    ("📰", "Fil d'actualité", "Posts multi-médias,\nlikes, commentaires"),
    ("📖", "Stories",         "Stories 24h avec\nréactions emoji"),
    ("💬", "Messagerie",      "DM + groupes,\nimages & audio"),
    ("📡", "Canaux campus",   "Canaux thématiques\npar filière"),
    ("📚", "Bibliothèque",    "Ressources partagées\npar filière"),
    ("🔔", "Notifications",   "Activité en temps\nréel"),
    ("👮", "Admin/Modération","Gestion utilisateurs\net contenu"),
    ("🏛️", "Multi-tenant",    "Données isolées\npar université"),
]

for i, (icon, title, desc) in enumerate(features):
    col = i % 4
    row = i // 4
    x = 0.35 + col * 3.25
    y = 2.75 + row * 2.2
    card(s, x, y, 3.05, 1.95, bg=C_PRIMARY)
    add_text(s, icon,  x + 0.15, y + 0.1, 0.7, 0.6, size=24, color=C_WHITE)
    add_text(s, title, x + 0.15, y + 0.7, 2.7, 0.5,
             size=14, bold=True, color=C_ACCENT)
    add_text(s, desc,  x + 0.15, y + 1.2, 2.7, 0.65,
             size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

footer(s, 6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Besoins Fonctionnels
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Besoins Fonctionnels",
             "Ce que le système doit faire")

left_items = [
    "Inscription / Connexion par email universitaire",
    "Vérification d'email & gestion du mot de passe",
    "Création de posts (texte + médias multiples)",
    "Liker, sauvegarder et commenter des posts",
    "Création et visualisation de stories (24h)",
    "Réactions emoji sur les stories",
    "Messagerie directe (1-à-1) et groupes",
    "Envoi d'images et d'audios dans les messages",
]
right_items = [
    "Abonnement à des utilisateurs et des canaux",
    "Profil personnalisable (bio, avatar, avatar-builder)",
    "Bibliothèque de ressources partagées par filière",
    "Système de notifications en temps réel",
    "Tableau de bord admin (modération, bannissement)",
    "Super-admin : gestion multi-universités",
    "Recherche et filtrage du fil d'actualité",
    "Changement de langue (AR / FR / EN)",
]

card(s, 0.3, 1.55, 6.3, 5.5, bg=C_WHITE)
card(s, 6.8, 1.55, 6.3, 5.5, bg=C_WHITE)
add_rect(s, 0.3, 1.55, 6.3, 0.07, fill=C_ACCENT)
add_rect(s, 6.8, 1.55, 6.3, 0.07, fill=C_GREEN)

add_text(s, "👤 Utilisateur / Étudiant", 0.5, 1.65, 5.9, 0.5,
         size=15, bold=True, color=C_PRIMARY)
add_text(s, "⚙️ Admin & Système", 7.0, 1.65, 5.9, 0.5,
         size=15, bold=True, color=C_PRIMARY)

bullet_box(s, left_items,  0.5, 2.2, 6.0, 4.5, text_size=13)
bullet_box(s, right_items, 7.0, 2.2, 6.0, 4.5, text_size=13, dot_color=C_GREEN)

footer(s, 7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Besoins Non Fonctionnels
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Besoins Non Fonctionnels",
             "Qualités transversales du système")

nfr = [
    ("🔒", "Sécurité",
     ["Auth Laravel Sanctum + CSRF", "Validation stricte des entrées", "Bannissement & middleware de rôles", "Mots de passe hachés (bcrypt)"]),
    ("⚡", "Performance",
     ["Eager loading (N+1 évité)", "Pagination du fil d'actualité", "Cache Inertia & compilation Vite", "Assets optimisés (chunking)"]),
    ("🌐", "Accessibilité",
     ["3 langues : AR / FR / EN", "Support RTL automatique", "Interface responsive (mobile-first)", "Focus mode (22h – 7h)"]),
    ("📦", "Maintenabilité",
     ["Architecture MVC Laravel", "Composants React réutilisables", "Migrations versionnées (24)", "Code séparé par domaine"]),
    ("🏫", "Multi-tenancy",
     ["Isolation par university_id", "Domaines email par campus", "Aliasing de domaines", "Admin par université"]),
    ("🔧", "Fiabilité",
     ["Gestion des erreurs (404/500)", "Notifications auto-purgées", "Stories auto-expirées (24h)", "Validation des fichiers uploadés"]),
]

for i, (icon, title, items) in enumerate(nfr):
    col = i % 3
    row = i // 2
    x = 0.3 + col * 4.35
    y = 1.6 + row * 2.7
    card(s, x, y, 4.1, 2.55, bg=C_WHITE)
    add_rect(s, x, y, 4.1, 0.5, fill=C_PRIMARY)
    add_text(s, f"{icon}  {title}", x + 0.15, y + 0.05, 3.8, 0.45,
             size=15, bold=True, color=C_WHITE)
    bullet_box(s, items, x + 0.1, y + 0.6, 3.9, 1.85,
               text_size=12, dot_color=C_ACCENT)

footer(s, 8)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Conception — Architecture Globale
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Conception — Architecture Globale",
             "Architecture MVC full-stack avec Inertia.js")

# Three-tier diagram
layers = [
    ("🖥️ Couche Présentation", C_ACCENT,
     ["React 18 + JSX", "Tailwind CSS 3.2", "Inertia.js (SPA sans API REST)", "Vite 7 (build)", "react-i18next"]),
    ("⚙️ Couche Métier", C_PRIMARY,
     ["Laravel 12 Controllers", "Eloquent ORM (14 modèles)", "Middleware (auth, rôles, locale)", "Form Requests (validation)", "Services & Helpers"]),
    ("🗄️ Couche Données", C_GREEN,
     ["SQLite (portable)", "24 migrations", "Stockage fichiers (storage/)", "Ziggy (routes JS)", "Laravel Sanctum"]),
]

for i, (label, color, items) in enumerate(layers):
    x = 0.3 + i * 4.35
    card(s, x, 1.6, 4.1, 5.5, bg=color)
    add_text(s, label, x + 0.15, 1.7, 3.8, 0.6,
             size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txb = s.shapes.add_textbox(Inches(x + 0.15), Inches(2.4), Inches(3.8), Inches(4.4))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(8)
        r1 = p.add_run(); r1.text = "●  "; r1.font.color.rgb = C_WHITE; r1.font.size = Pt(13)
        r2 = p.add_run(); r2.text = item; r2.font.color.rgb = C_WHITE; r2.font.size = Pt(13)

    if i < 2:
        add_text(s, "→", x + 4.15, 4.0, 0.2, 0.5,
                 size=26, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

footer(s, 9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Conception — Modèle de données
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Conception — Modèle de Données",
             "Les 14 entités principales et leurs relations")

entities = [
    # (label, x, y, color)
    ("University",    0.3,  1.65, C_ORANGE),
    ("User",          3.2,  1.65, C_PRIMARY),
    ("Post",          6.1,  1.65, C_ACCENT),
    ("PostItem",      9.0,  1.65, C_ACCENT),
    ("PostMedia",     11.1, 1.65, C_ACCENT),
    ("Story",         0.3,  3.55, C_GREEN),
    ("StoryView",     2.4,  3.55, C_GREEN),
    ("Channel",       4.8,  3.55, C_PRIMARY),
    ("Conversation",  7.2,  3.55, C_PRIMARY),
    ("Message",       9.6,  3.55, C_PRIMARY),
    ("Resource",      0.3,  5.35, C_ORANGE),
    ("PostComment",   3.2,  5.35, C_RED),
    ("Notification",  6.1,  5.35, C_GRAY),
    ("ContactMsg",    9.0,  5.35, C_GRAY),
]

for label, x, y, color in entities:
    card(s, x, y, 1.95, 0.6, bg=color)
    add_text(s, label, x + 0.05, y + 0.08, 1.85, 0.45,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Relations legend
legend = [
    (C_ORANGE,  "Gestion universitaire"),
    (C_PRIMARY, "Utilisateurs & Social"),
    (C_ACCENT,  "Contenu (Posts)"),
    (C_GREEN,   "Stories"),
    (C_RED,     "Commentaires"),
    (C_GRAY,    "Notifications & Contact"),
]
for i, (color, label) in enumerate(legend):
    x = 0.3 + i * 2.15
    add_rect(s, x, 6.55, 0.3, 0.3, fill=color)
    add_text(s, label, x + 0.38, 6.52, 1.7, 0.35, size=11, color=C_DARK_TEXT)

footer(s, 10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Conception — Diagramme de cas d'utilisation
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Conception — Cas d'Utilisation",
             "Principaux acteurs et leurs interactions avec le système")

actors = [
    ("👤 Étudiant", 0.3, [
        "S'inscrire / Se connecter",
        "Publier un post ou une story",
        "Liker / commenter / sauvegarder",
        "Envoyer des messages",
        "Suivre des utilisateurs / canaux",
        "Partager une ressource",
        "Consulter les notifications",
        "Modifier son profil",
    ]),
    ("👮 Admin Campus", 4.75, [
        "Modérer les posts",
        "Bannir / débannir un utilisateur",
        "Gérer les canaux campus",
        "Attribuer les rôles campus",
        "Gérer la bibliothèque",
        "Lire les messages de contact",
    ]),
    ("🏛️ Super-Admin", 9.2, [
        "Créer / gérer des universités",
        "Gérer les admins par campus",
        "Configurer les domaines email",
        "Aliaser des domaines",
        "Lire toutes les demandes",
    ]),
]

for label, x, items in actors:
    card(s, x, 1.55, 4.0, 5.55, bg=C_WHITE, border=C_ACCENT)
    add_rect(s, x, 1.55, 4.0, 0.55, fill=C_PRIMARY)
    add_text(s, label, x + 0.1, 1.6, 3.8, 0.45,
             size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    bullet_box(s, items, x + 0.1, 2.2, 3.8, 4.7, text_size=13)

footer(s, 11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Implémentation — Stack Technique
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Implémentation — Stack Technique",
             "Technologies choisies et justification")

techs = [
    ("Laravel 12", "PHP 8.2+",
     "Framework MVC robuste, Eloquent ORM, Sanctum pour l'auth, routing expressif, middleware puissant.",
     RGBColor(0xFF, 0x29, 0x1D)),
    ("React 18", "Frontend",
     "Composants réutilisables, hooks, Context API pour la gestion d'état global du feed.",
     RGBColor(0x61, 0xDA, 0xFB)),
    ("Inertia.js 2", "Glue fullstack",
     "Élimine le besoin d'une API REST séparée. Les props Laravel deviennent des props React directement.",
     RGBColor(0x92, 0x58, 0xFA)),
    ("Tailwind CSS 3", "Styling",
     "Utilitaires CSS, thème dark/light via CSS variables, responsive design mobile-first.",
     RGBColor(0x38, 0xBD, 0xF8)),
    ("SQLite", "Base de données",
     "Portable, zéro configuration serveur — idéal pour un MVP académique et le développement local.",
     RGBColor(0x00, 0x3B, 0x57)),
    ("Vite 7", "Build tool",
     "HMR ultra-rapide, code-splitting automatique, bundling optimisé pour la production.",
     RGBColor(0xBD, 0x34, 0xFE)),
]

for i, (name, tag, desc, color) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = 0.3 + col * 4.35
    y = 1.6 + row * 2.7
    card(s, x, y, 4.1, 2.5, bg=C_WHITE)
    add_rect(s, x, y, 0.08, 2.5, fill=color)
    add_text(s, name, x + 0.25, y + 0.1, 3.7, 0.5,
             size=17, bold=True, color=C_PRIMARY)
    add_text(s, tag,  x + 0.25, y + 0.55, 3.7, 0.35,
             size=12, color=color, bold=True)
    add_text(s, desc, x + 0.25, y + 0.9,  3.7, 1.45,
             size=12, color=C_GRAY)

footer(s, 12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Implémentation — Fonctionnalités clés (Feed & Posts)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Implémentation — Fil d'Actualité & Posts",
             "Cœur du réseau social — création, partage et interaction")

card(s, 0.3, 1.6, 8.0, 5.5, bg=C_WHITE)
add_text(s, "Fonctionnement du fil d'actualité", 0.5, 1.7, 7.6, 0.45,
         size=16, bold=True, color=C_PRIMARY)

feed_items = [
    "Fil agrégé : posts des utilisateurs suivis + canaux suivis + tous les posts du campus",
    "Tri chronologique avec pagination (infinite scroll / bouton)",
    "Filtre par canal campus via la barre de préférences",
    "Posts multi-blocs : texte + jusqu'à 10 médias (15 Mo chacun)",
    "Système PostItem + PostMedia pour les slides ordonnées",
    "Karma +1 à chaque publication (gamification)",
    "Filtrage automatique de mots interdits à la création",
    "Prévisualisation des 50 derniers likers avec avatars",
]
bullet_box(s, feed_items, 0.5, 2.25, 7.6, 4.5, text_size=14)

card(s, 8.55, 1.6, 4.45, 2.55, bg=C_PRIMARY)
add_text(s, "PostController", 8.7, 1.7, 4.1, 0.45,
         size=14, bold=True, color=C_ACCENT)
code_lines = [
    "store()        → créer post",
    "destroy()      → supprimer",
    "toggleLike()   → like/unlike",
    "toggleSave()   → sauvegarder",
    "storeComment() → commenter",
]
for j, line in enumerate(code_lines):
    add_text(s, line, 8.7, 2.2 + j * 0.36, 4.1, 0.35,
             size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

card(s, 8.55, 4.35, 4.45, 2.75, bg=C_LIGHT)
add_text(s, "DashboardController", 8.7, 4.45, 4.1, 0.45,
         size=14, bold=True, color=C_PRIMARY)
dash_lines = [
    "→ Agrège posts + stories",
    "→ Trending hashtags",
    "→ Utilisateurs en ligne",
    "→ Suggestions de suivi",
    "→ Stories du campus",
]
for j, line in enumerate(dash_lines):
    add_text(s, line, 8.7, 4.95 + j * 0.36, 4.1, 0.35,
             size=11, color=C_GRAY)

footer(s, 13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Implémentation — Messagerie & Stories
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Implémentation — Messagerie & Stories",
             "Communication en temps réel et contenu éphémère")

# Messaging
card(s, 0.3, 1.6, 6.1, 5.5, bg=C_WHITE)
add_rect(s, 0.3, 1.6, 6.1, 0.55, fill=C_PRIMARY)
add_text(s, "💬  Messagerie", 0.45, 1.65, 5.8, 0.45,
         size=16, bold=True, color=C_WHITE)

msg_items = [
    "Messages directs (1-à-1) et groupes",
    "user_one_id / user_two_id normalisés (min/max) → unicité",
    "Pivot conversation_user pour les groupes (N membres)",
    "Types de contenu : texte, image, audio",
    "read_at timestamp → accusés de lecture",
    "Création de groupe avec titre et liste de membres",
    "Filtrage des contacts par nom d'utilisateur",
]
bullet_box(s, msg_items, 0.5, 2.3, 5.8, 4.5, text_size=13)

# Stories
card(s, 6.7, 1.6, 6.3, 5.5, bg=C_WHITE)
add_rect(s, 6.7, 1.6, 6.3, 0.55, fill=C_GREEN)
add_text(s, "📖  Stories", 6.85, 1.65, 5.9, 0.45,
         size=16, bold=True, color=C_WHITE)

story_items = [
    "Durée de vie : 24 heures (expires_at)",
    "Médias ou texte court avec couleur de fond",
    "Vignettes circulaires dans la StoriesBar",
    "StoryViewerModal : carousel avec navigation",
    "Réactions emoji : 😍 🤣 😢 😡 👍",
    "Compteur de vues par story (table story_views)",
    "Auto-exclusion des propres stories vues",
]
bullet_box(s, story_items, 6.85, 2.3, 5.9, 4.5, text_size=13, dot_color=C_GREEN)

footer(s, 14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Implémentation — Multi-tenancy & Architecture éthique
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Implémentation — Multi-Tenancy & Éthique",
             "Isolation des données et conception responsable")

# Multi-tenancy
card(s, 0.3, 1.6, 6.1, 5.5, bg=C_WHITE)
add_rect(s, 0.3, 1.6, 6.1, 0.55, fill=C_ORANGE)
add_text(s, "🏫  Architecture Multi-Campus", 0.45, 1.65, 5.8, 0.45,
         size=15, bold=True, color=C_WHITE)

mt_items = [
    "Chaque université possède son propre espace isolé",
    "Inscription uniquement avec l'email du campus (domaine validé)",
    "university_id sur chaque table de données",
    "Aliasing de domaines : plusieurs emails → 1 campus",
    "Réassignation des données lors de la suppression d'un alias",
    "Admin dédié par campus (role = 'admin')",
    "Super-admin global (role = 'super_admin')",
]
bullet_box(s, mt_items, 0.5, 2.3, 5.8, 4.5, text_size=13, dot_color=C_ORANGE)

# Ethical design
card(s, 6.7, 1.6, 6.3, 5.5, bg=C_WHITE)
add_rect(s, 6.7, 1.6, 6.3, 0.55, fill=C_ACCENT)
add_text(s, "🛡️  Conception Éthique", 6.85, 1.65, 5.9, 0.45,
         size=15, bold=True, color=C_WHITE)

eth_items = [
    "Mode Focus (22h–7h) : limite les distractions nocturnes",
    "Système de karma : +1 par post pour encourager la participation",
    "Filtrage de mots interdits à la publication",
    "Formulaire de contact anonyme possible",
    "Vérification d'email avant accès au réseau",
    "Middleware EnsureUserIsNotBanned (déconnexion immédiate)",
    "Suppression de compte + données (RGPD-friendly)",
    "Notifications auto-supprimées après 1 jour (lues)",
]
bullet_box(s, eth_items, 6.85, 2.3, 5.9, 4.5, text_size=13)

footer(s, 15)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Implémentation — Admin & Modération
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Implémentation — Administration & Modération",
             "Outils de gestion pour les administrateurs")

admin_sections = [
    ("👮 Admin Campus", C_PRIMARY, [
        "Tableau de bord avec stats en temps réel",
        "Supprimer n'importe quel post du campus",
        "Bannir / débannir des utilisateurs",
        "Modifier le rôle campus (étudiant/enseignant/staff)",
        "Créer et supprimer des canaux thématiques",
        "Gérer les ressources de la bibliothèque",
        "Lire les messages de contact du campus",
    ]),
    ("🏛️ Super-Admin", C_ORANGE, [
        "Créer / éditer / supprimer des universités",
        "Configurer le domaine email par campus",
        "Ajouter des alias de domaines secondaires",
        "Créer et gérer les comptes admin",
        "Voir toutes les demandes de contact",
        "Marquer les messages comme lus",
    ]),
    ("🛡️ Middleware & Sécurité", C_GREEN, [
        "IsAdmin → protège les routes /admin/*",
        "IsSuperAdmin → protège /super-admin/*",
        "EnsureUserIsNotBanned → logout immédiat",
        "HandleInertiaRequests → partage global",
        "SetLocale → locale par session",
        "Throttle sur le formulaire de contact",
    ]),
]

for i, (title, color, items) in enumerate(admin_sections):
    x = 0.3 + i * 4.35
    card(s, x, 1.6, 4.1, 5.5, bg=C_WHITE)
    add_rect(s, x, 1.6, 4.1, 0.55, fill=color)
    add_text(s, title, x + 0.15, 1.65, 3.8, 0.45,
             size=14, bold=True, color=C_WHITE)
    bullet_box(s, items, x + 0.15, 2.25, 3.8, 4.6, text_size=12, dot_color=color)

footer(s, 16)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Interface utilisateur
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Implémentation — Interface Utilisateur",
             "26 pages React, 41 composants, 2 thèmes")

pages = [
    ("🏠", "Welcome",       "Page d'accueil publique"),
    ("📰", "Dashboard",     "Fil d'actualité principal"),
    ("📖", "Stories",       "Stories 24h"),
    ("💬", "Messages",      "Messagerie DM & groupes"),
    ("📡", "Channels",      "Canaux campus"),
    ("📚", "Library",       "Bibliothèque de ressources"),
    ("🔔", "Notifications", "Fil d'activité"),
    ("👤", "Profile",       "Profil public & édition"),
    ("🔖", "Saved Posts",   "Posts sauvegardés"),
    ("👮", "Admin",         "Dashboard modération"),
    ("🏛️", "Super-Admin",   "Gestion universités"),
    ("💡", "Suggestions",   "Suggestions d'abonnement"),
]

for i, (icon, name, desc) in enumerate(pages):
    col = i % 4
    row = i // 4
    x = 0.3 + col * 3.25
    y = 1.6 + row * 1.85
    card(s, x, y, 3.05, 1.65, bg=C_WHITE)
    add_text(s, icon, x + 0.15, y + 0.1,  0.6, 0.6, size=22)
    add_text(s, name, x + 0.8, y + 0.1,  2.1, 0.5,
             size=14, bold=True, color=C_PRIMARY)
    add_text(s, desc, x + 0.15, y + 0.9,  2.8, 0.6,
             size=11, color=C_GRAY)

# UI features note
card(s, 0.3, 7.0, 12.73, 0.35, bg=C_PRIMARY)
add_text(s, "✦  Thème Dark / Light  ✦  Mode compact sidebar  ✦  RTL automatique pour l'arabe  ✦  Toasts de notifications",
         0.5, 7.02, 12.3, 0.3, size=12, color=C_WHITE, align=PP_ALIGN.CENTER)

footer(s, 17)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_LIGHT)
slide_header(s, "Conclusion",
             "Bilan du projet UniConnect")

# Achieved
card(s, 0.3, 1.6, 5.9, 5.5, bg=C_WHITE, border=C_GREEN)
add_rect(s, 0.3, 1.6, 5.9, 0.55, fill=C_GREEN)
add_text(s, "✅  Réalisations", 0.45, 1.65, 5.6, 0.45,
         size=16, bold=True, color=C_WHITE)

achieved = [
    "Réseau social académique complet et fonctionnel",
    "8 fonctionnalités sociales principales",
    "Architecture multi-tenant (N universités)",
    "14 modèles, 19 contrôleurs, 136 routes",
    "Interface trilingue AR/FR/EN avec RTL",
    "Conception éthique (focus, karma, filtrage)",
    "Panneau admin + super-admin complet",
    "Code maintenable : MVC + composants React",
]
bullet_box(s, achieved, 0.5, 2.3, 5.6, 4.5, text_size=13, dot_color=C_GREEN)

# Perspectives
card(s, 6.5, 1.6, 6.5, 5.5, bg=C_WHITE, border=C_ACCENT)
add_rect(s, 6.5, 1.6, 6.5, 0.55, fill=C_ACCENT)
add_text(s, "🚀  Perspectives & Améliorations", 6.65, 1.65, 6.2, 0.45,
         size=16, bold=True, color=C_WHITE)

perspectives = [
    "Temps réel avec Laravel Broadcasting / WebSockets",
    "Application mobile (React Native)",
    "Notifications push (FCM / Web Push)",
    "IA : recommandations de contenu personnalisées",
    "Chatbot académique intégré (aide étudiante)",
    "Tests unitaires et d'intégration complets",
    "Déploiement cloud (Docker + CI/CD)",
    "API publique pour intégrations tierces",
]
bullet_box(s, perspectives, 6.65, 2.3, 6.2, 4.5, text_size=13)

footer(s, 18)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Merci / Questions
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=C_PRIMARY)
add_rect(s, 0, 3.3, 13.33, 0.08, fill=C_ACCENT)

add_text(s, "Merci pour votre attention !",
         0.5, 0.8, 12.33, 1.3,
         size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(s, "🎓  UniConnect — Réseau Social Académique",
         0.5, 2.1, 12.33, 0.7,
         size=22, color=C_ACCENT, align=PP_ALIGN.CENTER)

add_text(s, "Des questions ?",
         0.5, 3.6, 12.33, 0.9,
         size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Summary stats
stats = [
    ("14", "Modèles"),
    ("19+", "Contrôleurs"),
    ("136", "Routes"),
    ("26", "Pages React"),
    ("41", "Composants"),
    ("3", "Langues"),
]
for i, (val, label) in enumerate(stats):
    x = 0.8 + i * 2.0
    add_rect(s, x, 4.75, 1.7, 1.4, fill=RGBColor(0x16, 0x27, 0x4E))
    add_text(s, val,   x + 0.05, 4.82, 1.6, 0.7,
             size=32, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, label, x + 0.05, 5.5,  1.6, 0.45,
             size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

footer(s, 19)

# ─── Save ────────────────────────────────────────────────────────────────────
out = r"c:\Users\elarr\2026\uniconnect\UniConnect_Presentation.pptx"
prs.save(out)
print(f"OK - Fichier genere : {out}")
